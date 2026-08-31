import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def build_presentation():
    prs = pptx.Presentation('C:/fra-atlas/SIH2026-IDEA-Presentation-Format.pptx')
    
    # -------------------------------------------------------------
    # SLIDE 1: Title Page
    # -------------------------------------------------------------
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                text = p.text
                if 'Problem Statement ID' in text:
                    p.text = "Problem Statement ID : SIH12508"
                    p.font.size = Pt(13)
                    p.font.bold = True
                elif 'Problem Statement Title' in text:
                    p.text = "Problem Statement Title : Real-Time Map Tracking & Document Digitization for Forest Rights Claims"
                    p.font.size = Pt(13)
                    p.font.bold = True
                elif 'Theme-' in text or text.strip() == 'Theme-':
                    p.text = "Theme : Smart Automation / Clean & Green Governance"
                    p.font.size = Pt(13)
                    p.font.bold = True
                elif 'PS Category' in text:
                    p.text = "PS Category : Software"
                    p.font.size = Pt(13)
                    p.font.bold = True
                elif 'Team ID-' in text or text.strip() == 'Team ID-':
                    p.text = "Team ID : SIH2026-FRA-ATLAS"
                    p.font.size = Pt(13)
                    p.font.bold = True
                elif 'Team Name' in text:
                    p.text = "Team Name : FRA Atlas Innovators"
                    p.font.size = Pt(13)
                    p.font.bold = True

    # -------------------------------------------------------------
    # SLIDE 2: Idea Title & Proposed Solution
    # -------------------------------------------------------------
    slide2 = prs.slides[1]
    # Set Idea Title
    for shape in slide2.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'IDEA TITLE' in p.text or p.text.strip() == 'IDEA TITLE':
                    p.text = "FRA Atlas: AI-Powered WebGIS, Document OCR & Blockchain Title Platform"
                    p.font.size = Pt(16)
                    p.font.bold = True
                elif 'Your Team Name' in p.text:
                    p.text = "FRA Atlas Innovators"
                elif 'Proposed Solution' in p.text:
                    # Clear default placeholder text
                    shape.text_frame.text = ""
                    tf = shape.text_frame
                    
                    # 1. Proposed Solution
                    p1 = tf.paragraphs[0]
                    p1.text = "1. Proposed Solution (FRA Atlas Prototype):"
                    p1.font.bold = True
                    p1.font.size = Pt(13)
                    p1.font.color.rgb = RGBColor(22, 101, 52)
                    
                    p1_sub = tf.add_paragraph()
                    p1_sub.text = "• FRA Atlas is a centralized WebGIS & AI platform that solves 18+ years of delayed Forest Rights Act (2006) claims for 2+ crore tribal citizens.\n• Delivers real-time 36-state choropleth tracking, automatic OCR digitization, machine learning fraud surveillance, and immutable blockchain titles."
                    p1_sub.font.size = Pt(11)
                    
                    # 2. How it addresses problem
                    p2 = tf.add_paragraph()
                    p2.text = "2. How It Addresses Critical Bottlenecks:"
                    p2.font.bold = True
                    p2.font.size = Pt(13)
                    p2.font.color.rgb = RGBColor(22, 101, 52)
                    
                    p2_sub = tf.add_paragraph()
                    p2_sub.text = "• Eliminates 20+ Lakh Pending Backlog: Digitize legacy paper claims into structured spatial records in under 2 seconds.\n• Transparency on Rejections: Reveals specific Section 12A rejection codes for immediate Gram Sabha legal appeal.\n• Transparent Governance: Replaces delayed manual reports with a live nationwide administrative Command Center."
                    p2_sub.font.size = Pt(11)

                    # 3. Innovation & Uniqueness
                    p3 = tf.add_paragraph()
                    p3.text = "3. Innovation & Uniqueness:"
                    p3.font.bold = True
                    p3.font.size = Pt(13)
                    p3.font.color.rgb = RGBColor(22, 101, 52)
                    
                    p3_sub = tf.add_paragraph()
                    p3_sub.text = "• Multi-Layer Integration: Combines MapLibre WebGIS + PaddleOCR Vision + Isolation Forest ML + SHA-256 Merkle Ledger.\n• 100% Free & Open-Source Stack: Zero recurring licensing fee (MIT / BSD-3 tools), ready for nationwide deployment."
                    p3_sub.font.size = Pt(11)

    # -------------------------------------------------------------
    # SLIDE 3: Technical Approach
    # -------------------------------------------------------------
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'Your Team Name' in p.text:
                    p.text = "FRA Atlas Innovators"
                elif 'Technologies to be used' in p.text:
                    shape.text_frame.text = ""
                    tf = shape.text_frame
                    
                    # Section 1
                    p1 = tf.paragraphs[0]
                    p1.text = "1. Technology Stack Architecture (Zero-Cost Free Stack):"
                    p1.font.bold = True
                    p1.font.size = Pt(13)
                    p1.font.color.rgb = RGBColor(22, 101, 52)
                    
                    p1_sub = tf.add_paragraph()
                    p1_sub.text = "• Frontend: Next.js 16 + React 19 + Tailwind CSS + MapLibre GL JS (WebGL 60 FPS hardware accelerated).\n• Backend: FastAPI (Python 3.11) + Uvicorn Async Server + SQLite/PostgreSQL Spatial DB.\n• AI / ML Engine: PaddleOCR (Multi-angle form extraction) + Scikit-Learn (Isolation Forest Outlier Detection).\n• Security & Integrity: SHA-256 Cryptographic Hash Chain Ledger (Deterministic Merkle proof)."
                    p1_sub.font.size = Pt(11)
                    
                    # Section 2
                    p2 = tf.add_paragraph()
                    p2.text = "2. Implementation Workflow & Methodology:"
                    p2.font.bold = True
                    p2.font.size = Pt(13)
                    p2.font.color.rgb = RGBColor(22, 101, 52)
                    
                    p2_sub = tf.add_paragraph()
                    p2_sub.text = "• Data Ingestion: Scrapes live MoTA statistics (forestrights.nic.in) + vectorizes 36 state/district boundaries.\n• OCR Vision Pipeline: Scans Form-A/Form-B → extracts claimant, village, tribe, area → auto-zooms WebGIS.\n• AI Surveillance: Evaluates processing velocity (<4 days) & rejection spikes to flag corruption anomalies.\n• Title Minting: Computes deterministic hash for Gram Sabha resolution + coordinates → prevents land grabbing."
                    p2_sub.font.size = Pt(11)

    # -------------------------------------------------------------
    # SLIDE 4: Feasibility and Viability
    # -------------------------------------------------------------
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'Your Team Name' in p.text:
                    p.text = "FRA Atlas Innovators"
                elif 'Analysis of the feasibility' in p.text:
                    shape.text_frame.text = ""
                    tf = shape.text_frame
                    
                    # Section 1
                    p1 = tf.paragraphs[0]
                    p1.text = "1. Feasibility Analysis (Technical & Operational):"
                    p1.font.bold = True
                    p1.font.size = Pt(13)
                    p1.font.color.rgb = RGBColor(22, 101, 52)
                    
                    p1_sub = tf.add_paragraph()
                    p1_sub.text = "• High Technical Feasibility: Built entirely on modular microservices with verified working prototype.\n• Low Compute Footprint: CPU-optimized PaddleOCR and lightweight MapLibre vector tiles run on standard government servers.\n• Operational Scalability: Seamlessly ingests all 54+ Lakh claims across 700+ districts without third-party API costs."
                    p1_sub.font.size = Pt(11)
                    
                    # Section 2
                    p2 = tf.add_paragraph()
                    p2.text = "2. Potential Challenges & Mitigation Strategies:"
                    p2.font.bold = True
                    p2.font.size = Pt(13)
                    p2.font.color.rgb = RGBColor(22, 101, 52)
                    
                    p2_sub = tf.add_paragraph()
                    p2_sub.text = "• Low Bandwidth in Remote Forest Villages → Offline PWA caching + lightweight GeoJSON vector tiles.\n• Poor Quality Scanned Paper Forms → Pre-processing binarization + confidence threshold scoring in PaddleOCR.\n• Resistance to Digital Transition → Intuitive multi-stakeholder interface + automated sample demo flows.\n• Cadastral Map Alignment → Standardized WGS84 coordinates mapped to Survey of India boundaries."
                    p2_sub.font.size = Pt(11)

    # -------------------------------------------------------------
    # SLIDE 5: Impact and Benefits
    # -------------------------------------------------------------
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'Your Team Name' in p.text:
                    p.text = "FRA Atlas Innovators"
                elif 'Potential impact on the target' in p.text:
                    shape.text_frame.text = ""
                    tf = shape.text_frame
                    
                    # Section 1
                    p1 = tf.paragraphs[0]
                    p1.text = "1. Measurable Social & Economic Impact:"
                    p1.font.bold = True
                    p1.font.size = Pt(13)
                    p1.font.color.rgb = RGBColor(22, 101, 52)
                    
                    p1_sub = tf.add_paragraph()
                    p1_sub.text = "• 2+ Crore Tribal Citizens: Direct visibility into ancestral land titles, preventing illegal eviction and disputes.\n• Welfare Entitlement Unlocking: Verified title holders gain instant linkage to PM-KISAN, PMAY-G, and MGNREGA.\n• 85% Reduction in Processing Latency: Transition from physical paper bureaucracy to 2-second AI digitization.\n• Protection of 2.5+ Lakh Gram Sabhas: Empowers grassroots village councils with verifiable spatial evidence."
                    p1_sub.font.size = Pt(11)
                    
                    # Section 2
                    p2 = tf.add_paragraph()
                    p2.text = "2. Environmental & UN Sustainable Development Goals (SDG) Alignment:"
                    p2.font.bold = True
                    p2.font.size = Pt(13)
                    p2.font.color.rgb = RGBColor(22, 101, 52)
                    
                    p2_sub = tf.add_paragraph()
                    p2_sub.text = "• SDG 1 (No Poverty) & SDG 10 (Reduced Inequalities): Redresses historical injustice for indigenous tribes.\n• SDG 15 (Life on Land): Tribal community stewardship proven to achieve lowest deforestation rates across India.\n• SDG 16 (Peace, Justice & Strong Institutions): Tamper-proof title ledger builds accountable governance."
                    p2_sub.font.size = Pt(11)

    # -------------------------------------------------------------
    # SLIDE 6: Research and References
    # -------------------------------------------------------------
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'Your Team Name' in p.text:
                    p.text = "FRA Atlas Innovators"
                elif 'Details / Links of the reference' in p.text:
                    shape.text_frame.text = ""
                    tf = shape.text_frame
                    
                    # Section 1
                    p1 = tf.paragraphs[0]
                    p1.text = "1. Government Data Sources & Statutory References:"
                    p1.font.bold = True
                    p1.font.size = Pt(13)
                    p1.font.color.rgb = RGBColor(22, 101, 52)
                    
                    p1_sub = tf.add_paragraph()
                    p1_sub.text = "• Ministry of Tribal Affairs (MoTA) Official Portal: https://forestrights.nic.in (Monthly FRA Progress Reports).\n• Open Government Data (OGD) Platform India: https://data.gov.in (National FRA State Statistics Dataset).\n• Scheduled Tribes and Other Traditional Forest Dwellers (Recognition of Forest Rights) Act, 2006 (Act No. 2 of 2007).\n• Forest Survey of India (FSI): India State of Forest Report (ISFR) Geospatial Baselines."
                    p1_sub.font.size = Pt(11)
                    
                    # Section 2
                    p2 = tf.add_paragraph()
                    p2.text = "2. Academic Research & Technical Literature:"
                    p2.font.bold = True
                    p2.font.size = Pt(13)
                    p2.font.color.rgb = RGBColor(22, 101, 52)
                    
                    p2_sub = tf.add_paragraph()
                    p2_sub.text = "• Liu, F. T., Ting, K. M., & Zhou, Z. H. (2008). Isolation Forest. IEEE International Conference on Data Mining (ICDM).\n• Du, Y., et al. (2020). PaddleOCR: An Ultra Lightweight OCR System. Baidu Inc. Research.\n• MapLibre GL Technical Specification & WebGL Raster/Vector Tile Rendering Architecture."
                    p2_sub.font.size = Pt(11)

    # Note: Slide 7 is the instructions slide, kept as-is or deleted per user preference.
    
    prs.save('C:/fra-atlas/SIH2026_FRA_Atlas_Submission.pptx')
    print("Successfully generated C:/fra-atlas/SIH2026_FRA_Atlas_Submission.pptx")

if __name__ == "__main__":
    build_presentation()
