import pptx

prs = pptx.Presentation('C:/fra-atlas/SIH2026-IDEA-Presentation-Format.pptx')
for i, slide in enumerate(prs.slides):
    print(f"=== SLIDE {i+1} ===")
    for s_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                if p.text.strip():
                    print(f"  Shape {s_idx}, Para {p_idx}: '{p.text}'")
